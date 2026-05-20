#!/usr/bin/env python3
"""
Convert pptx/pdf/docx/xlsx/xls to markdown.
Usage: python3 convert_to_md.py <file> [module] [topic] [lecturer]

Improvements over v1:
- PDFs use pymupdf (much better layout/equation preservation than pdfminer)
- PPTXs extract embedded images to an assets/ subfolder
- All files get YAML frontmatter automatically
"""

import sys
import os

# ── Converters ────────────────────────────────────────────────────────────────

def convert_pptx(path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    file_stem = os.path.splitext(os.path.basename(path))[0]
    # Each source file gets its own subfolder to avoid name collisions
    assets_dir = os.path.join(os.path.dirname(path), "assets", file_stem)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        images = []

        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    img_name = f"slide{i}-img{shape.shape_id}.{ext}"
                    os.makedirs(assets_dir, exist_ok=True)
                    img_path = os.path.join(assets_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    images.append(f"![[assets/{file_stem}/{img_name}]]")
                except Exception:
                    pass

        parts = []
        if texts:
            parts.append("\n\n".join(texts))
        if images:
            parts.append("\n".join(images))

        if parts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(parts))

    return "\n\n---\n\n".join(slides)


def convert_pdf(path):
    try:
        import fitz  # pymupdf
        file_stem = os.path.splitext(os.path.basename(path))[0]
        assets_dir = os.path.join(os.path.dirname(path), "assets", file_stem)
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            images = []
            # Extract embedded images
            for img_index, img in enumerate(page.get_images(), 1):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    # Skip tiny images (icons, decorations) — min 5KB
                    if len(image_bytes) < 5120:
                        continue
                    img_name = f"page{i}-img{img_index}.{image_ext}"
                    os.makedirs(assets_dir, exist_ok=True)
                    with open(os.path.join(assets_dir, img_name), "wb") as f:
                        f.write(image_bytes)
                    images.append(f"![[assets/{file_stem}/{img_name}]]")
                except Exception:
                    pass
            parts = []
            if text:
                parts.append(text)
            if images:
                parts.append("\n".join(images))
            if parts:
                pages.append(f"## Page {i}\n\n" + "\n\n".join(parts))
        doc.close()
        return "\n\n---\n\n".join(pages)
    except Exception:
        # Fallback to pdfminer if pymupdf fails
        from pdfminer.high_level import extract_text
        return extract_text(path)


def convert_docx(path):
    from docx import Document

    OMML = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
    WORD = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    def omml_to_text(elem):
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        ns  = elem.tag.split('}')[0].lstrip('{') if '}' in elem.tag else ''
        if ns == WORD:
            return ''
        if tag == 't':
            return elem.text or ''
        elif tag == 'r':
            return ''.join((c.text or '') for c in elem if c.tag.split('}')[-1] == 't')
        elif tag == 'sSup':
            base = sup = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'e':    base = omml_to_text(c)
                elif ct == 'sup': sup = omml_to_text(c)
            return f'{base}^{sup}' if len(sup) <= 2 else f'{base}^({sup})'
        elif tag == 'sSub':
            base = sub = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'e':    base = omml_to_text(c)
                elif ct == 'sub': sub = omml_to_text(c)
            return f'{base}_{sub}' if len(sub) <= 2 else f'{base}_({sub})'
        elif tag == 'sSubSup':
            base = sub = sup = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'e':     base = omml_to_text(c)
                elif ct == 'sub': sub  = omml_to_text(c)
                elif ct == 'sup': sup  = omml_to_text(c)
            return f'{base}_({sub})^({sup})'
        elif tag == 'f':
            num = den = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'num': num = omml_to_text(c)
                elif ct == 'den': den = omml_to_text(c)
            return f'({num})/({den})'
        elif tag == 'rad':
            deg = content = ''; hidden = False
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'radPr':
                    for gc in c:
                        if gc.tag.split('}')[-1] == 'degHide':
                            hidden = gc.get(f'{{{OMML}}}val','1') not in ('0','false','off')
                elif ct == 'deg': deg = omml_to_text(c)
                elif ct == 'e':   content = omml_to_text(c)
            return f'√({content})' if hidden or not deg.strip() else f'{deg}√({content})'
        elif tag == 'd':
            beg, end = '(', ')'; contents = []
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'dPr':
                    for gc in c:
                        gct = gc.tag.split('}')[-1]
                        if gct == 'begChr': beg = gc.get(f'{{{OMML}}}val','(') or ''
                        elif gct == 'endChr': end = gc.get(f'{{{OMML}}}val',')') or ''
                elif ct == 'e': contents.append(omml_to_text(c))
            return f'{beg}{", ".join(contents)}{end}'
        elif tag == 'nary':
            op = '∫'; sub = sup = content = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'naryPr':
                    for gc in c:
                        if gc.tag.split('}')[-1] == 'chr':
                            op = gc.get(f'{{{OMML}}}val','∫') or '∫'
                elif ct == 'sub': sub  = omml_to_text(c)
                elif ct == 'sup': sup  = omml_to_text(c)
                elif ct == 'e':   content = omml_to_text(c)
            return f'{op}_{{{sub}}}^{{{sup}}} {content}' if (sub or sup) else f'{op} {content}'
        elif tag == 'func':
            fname = content = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'fName': fname = omml_to_text(c)
                elif ct == 'e':   content = omml_to_text(c)
            return f'{fname}({content})'
        elif tag == 'bar':
            content = next((omml_to_text(c) for c in elem if c.tag.split('}')[-1]=='e'),'')
            return f'{content}̅'
        elif tag == 'acc':
            return next((omml_to_text(c) for c in elem if c.tag.split('}')[-1]=='e'),'')
        elif tag == 'limLow':
            base = lim = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'e': base = omml_to_text(c)
                elif ct == 'lim': lim = omml_to_text(c)
            return f'{base}_{{{lim}}}'
        elif tag == 'limUpp':
            base = lim = ''
            for c in elem:
                ct = c.tag.split('}')[-1]
                if ct == 'e': base = omml_to_text(c)
                elif ct == 'lim': lim = omml_to_text(c)
            return f'{base}^{{{lim}}}'
        elif tag == 'm':
            rows = []
            for c in elem:
                if c.tag.split('}')[-1] == 'mr':
                    cells = [omml_to_text(gc) for gc in c if gc.tag.split('}')[-1]=='e']
                    rows.append('  '.join(cells))
            return '[' + ' | '.join(rows) + ']'
        elif tag == 'eqArr':
            return '\n'.join(omml_to_text(c) for c in elem if c.tag.split('}')[-1]=='e')
        elif tag in ('rPr','dPr','sSupPr','sSubPr','radPr','naryPr','fPr','barPr',
                     'accPr','ctrlPr','mPr','mrPr','sty','limLoc','degHide',
                     'sSubSupPr','groupChrPr','borderBoxPr','boxPr','eqArrPr','pPr'):
            return ''
        else:
            return ''.join(omml_to_text(c) for c in elem)

    def para_text_with_math(para_elem):
        parts = []
        for child in para_elem:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            ns  = child.tag.split('}')[0].lstrip('{') if '}' in child.tag else ''
            if ns == OMML and tag in ('oMath','oMathPara'):
                math = omml_to_text(child).strip()
                if math:
                    parts.append(f'`{math}`')
            elif ns == WORD and tag == 'r':
                for gc in child:
                    if gc.tag.split('}')[-1] == 't':
                        parts.append(gc.text or '')
            else:
                sub = para_text_with_math(child)
                if sub:
                    parts.append(sub)
        return ''.join(parts)

    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para_text_with_math(para._element).strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            lines.append(f"# {text}")
        elif "Heading 2" in style:
            lines.append(f"## {text}")
        elif "Heading 3" in style:
            lines.append(f"### {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def convert_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    output = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        output.append(f"## Sheet: {sheet_name}\n")
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("| " + " | ".join(cells) + " |")
        if rows:
            separator = "| " + " | ".join(["---"] * len(rows[0].split("|")[1:-1])) + " |"
            output.append(rows[0])
            output.append(separator)
            output.extend(rows[1:])
        output.append("")
    return "\n".join(output)


def convert_xls(path):
    import xlrd
    wb = xlrd.open_workbook(path)
    output = []
    for sheet in wb.sheets():
        output.append(f"## Sheet: {sheet.name}\n")
        rows = []
        for rx in range(sheet.nrows):
            cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
            if any(c.strip() for c in cells):
                rows.append("| " + " | ".join(cells) + " |")
        if rows:
            separator = "| " + " | ".join(["---"] * sheet.ncols) + " |"
            output.append(rows[0])
            output.append(separator)
            output.extend(rows[1:])
        output.append("")
    return "\n".join(output)


# ── Frontmatter ───────────────────────────────────────────────────────────────

TOPIC_TAGS = {
    "thermodynamics":                          ["thermodynamics", "thermofluids"],
    "fluid-mechanics":                         ["fluid-mechanics", "thermofluids"],
    "mechanics/statics":                       ["statics", "mechanics"],
    "mechanics/solid-mechanics":               ["solid-mechanics", "mechanics"],
    "mechanics/dynamics":                      ["dynamics", "mechanics"],
    "mechanics":                               ["mechanics"],
    "manufacturing-materials/manufacturing":   ["manufacturing", "manufacturing-materials"],
    "manufacturing-materials/materials":       ["materials", "manufacturing-materials"],
    "manufacturing-materials":                 ["manufacturing-materials"],
}

def topic_tags(topic):
    t = topic.lower()
    for key, tags in TOPIC_TAGS.items():
        if key in t:
            return tags
    return []

def build_frontmatter(path, module, topic, lecturer):
    basename = os.path.splitext(os.path.basename(path))[0]
    tags = topic_tags(topic)
    tags_str = "[" + ", ".join(tags) + "]" if tags else "[]"
    return (
        f"---\n"
        f"module: {module}\n"
        f"topic: {topic}\n"
        f"lecturer: {lecturer}\n"
        f"module_index: \"[[/_modules/{module}]]\"\n"
        f"source_file: \"{basename}\"\n"
        f"tags: {tags_str}\n"
        f"---\n\n"
    )


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 convert_to_md.py <file> [module] [topic] [lecturer]")
        sys.exit(1)

    path     = sys.argv[1]
    module   = sys.argv[2] if len(sys.argv) > 2 else "UNKNOWN"
    topic    = sys.argv[3] if len(sys.argv) > 3 else "unknown"
    lecturer = sys.argv[4] if len(sys.argv) > 4 else "Unknown"

    if not os.path.exists(path):
        print(f"File not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    converters = {
        ".pptx": convert_pptx,
        ".pdf":  convert_pdf,
        ".docx": convert_docx,
        ".xlsx": convert_xlsx,
        ".xls":  convert_xls,
    }

    if ext not in converters:
        print(f"Unsupported format: {ext}", file=sys.stderr)
        sys.exit(1)

    content      = converters[ext](path)
    frontmatter  = build_frontmatter(path, module, topic, lecturer)
    out_path     = os.path.splitext(path)[0] + ".md"

    with open(out_path, "w", encoding="utf-8") as f:
        f.write(frontmatter + content)

    print(out_path)


if __name__ == "__main__":
    main()
